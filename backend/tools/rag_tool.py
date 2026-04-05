"""
backend/tools/rag_tool.py

RAG pipeline backed by Supabase Storage + pgvector.

Storage layout
--------------
  seis-files/{folder_name}/{filename}

DB tables used
--------------
  folders  (id, name)
  files    (id, folder_id, name, storage_path, size_bytes)
  embeddings (id, folder_id, file_id, chunk_text, embedding vector(384))

pgvector similarity is done via the match_chunks() SQL function
that must exist in Supabase (already created per user confirmation).

SQL function signature (cosine distance):
  match_chunks(query_embedding vector(384),
               target_folder_id uuid,
               match_count      int)
  RETURNS TABLE (chunk_text text, similarity float)

File reader functions, chunk_text(), and run() response shape are unchanged.
"""

import io
import logging
import os

from backend.services.embeddings import embed_texts
from backend.services.llm import generate_response
from backend.services.supabase_client import supabase, BUCKET, execute_with_retry

from pypdf import PdfReader
from docx import Document
from pptx import Presentation

from fastapi import HTTPException

logger = logging.getLogger(__name__)


# -------------------- FILE READERS (UNCHANGED) --------------------

def read_txt(content: bytes) -> str:
    return content.decode("utf-8")


def read_py(content: bytes) -> str:
    try:
        return content.decode("utf-8")
    except UnicodeDecodeError:
        return content.decode("latin-1")


def read_code(content: bytes) -> str:
    """Generic reader for .c .cpp .h .hpp .java .asm .s"""
    try:
        return content.decode("utf-8")
    except UnicodeDecodeError:
        return content.decode("latin-1")


def read_pdf(content: bytes) -> str:
    reader = PdfReader(io.BytesIO(content))
    text = ""
    for page in reader.pages:
        text += page.extract_text() or ""
    return text


def read_docx(content: bytes) -> str:
    doc = Document(io.BytesIO(content))
    return "\n".join([para.text for para in doc.paragraphs])


def read_pptx(content: bytes) -> str:
    presentation = Presentation(io.BytesIO(content))
    text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text.append(shape.text)
    return "\n".join(text)


def _read_file_content(filename: str, content: bytes) -> str:
    """Dispatch to the right reader based on file extension."""
    fl = filename.lower()
    if fl.endswith(".txt"):
        return read_txt(content)
    if fl.endswith(".py"):
        return read_py(content)
    if fl.endswith((".c", ".cpp", ".h", ".hpp", ".java", ".asm", ".s",
                    ".js", ".ts", ".jsx", ".tsx", ".md")):
        return read_code(content)
    if fl.endswith(".pdf"):
        return read_pdf(content)
    if fl.endswith(".docx"):
        return read_docx(content)
    if fl.endswith(".pptx"):
        return read_pptx(content)
    return ""   # unsupported — caller skips empty strings


# -------------------- CHUNKING (UNCHANGED) --------------------

def chunk_text(text, chunk_size=300, overlap=50):
    words = text.split()
    chunks = []
    for i in range(0, len(words), chunk_size - overlap):
        chunk = " ".join(words[i:i + chunk_size])
        chunks.append(chunk)
    return chunks


# -------------------- CONSTANTS --------------------

TOP_K = 3


# -------------------- EMBEDDING HELPERS --------------------

def _upsert_embeddings(file_id: str, chunks: list[str]):
    print(f"[_upsert_embeddings] file_id={file_id!r}  chunk_count={len(chunks)}", flush=True)
    try:
        supabase.table("embeddings").delete().eq("file_id", file_id).execute()
    except Exception:
        pass

    if not chunks:
        print(f"[_upsert_embeddings] WARNING: no chunks — skipping embed for file_id={file_id!r}", flush=True)
        return

    try:
        vectors = embed_texts(chunks)
    except Exception as e:
        import traceback
        print(f"[_upsert_embeddings] ERROR: embed_texts() FAILED for file_id={file_id!r}: {e}", flush=True)
        print(traceback.format_exc(), flush=True)
        raise

    print(f"[_upsert_embeddings] embedded {len(vectors)} chunks for file_id={file_id!r}", flush=True)

    rows = []
    for chunk, vec in zip(chunks, vectors):
        rows.append({
            "file_id":    file_id,
            "chunk_text": chunk,
            "embedding":  vec.tolist(),
        })

    batch = 100
    for start in range(0, len(rows), batch):
        try:
            supabase.table("embeddings").insert(rows[start:start + batch]).execute()
            print(f"[_upsert_embeddings] inserted batch {start}–{start + len(rows[start:start+batch])}", flush=True)
        except Exception as e:
            import traceback
            print(f"[_upsert_embeddings] ERROR: INSERT FAILED at batch {start}, file_id={file_id!r}: {e}", flush=True)
            print(traceback.format_exc(), flush=True)
            raise


# -------------------- INDEX A SINGLE FILE --------------------

def index_file(folder_id: str, file_id: str, storage_path: str, filename: str):
    print(f"[index_file] START  folder_id={folder_id!r}  file_id={file_id!r}  path={storage_path!r}", flush=True)
    try:
        raw = supabase.storage.from_(BUCKET).download(storage_path)
        print(f"[index_file] downloaded {len(raw)} bytes for {storage_path!r}", flush=True)
    except Exception as exc:
        import traceback
        print(f"[index_file] ERROR: Storage download FAILED for path={storage_path!r}: {exc}", flush=True)
        print(traceback.format_exc(), flush=True)
        raise HTTPException(status_code=503, detail=f"Storage unavailable: {exc}")

    text = _read_file_content(filename, raw)
    if not text.strip():
        print(f"[index_file] WARNING: extracted text is EMPTY for file={filename!r} — skipping embed", flush=True)
        return

    print(f"[index_file] extracted {len(text)} chars, creating chunks...", flush=True)
    chunks = chunk_text(text)
    print(f"[index_file] created {len(chunks)} chunks for {filename!r}", flush=True)
    _upsert_embeddings(file_id, chunks)
    print(f"[index_file] DONE for file_id={file_id!r}", flush=True)


# -------------------- INDEX WHOLE FOLDER --------------------

def index_folder(folder_id: str, folder_name: str):
    """
    (Re)index all files belonging to a folder.
    Fetches the file list from the `files` DB table, then calls index_file()
    for each entry.
    """
    try:
        resp = execute_with_retry(
            lambda sb: (
                sb.table("files")
                .select("id, name, storage_path")
                .eq("folder_id", folder_id)
                .execute()
            )
        )
    except Exception as exc:
        raise HTTPException(status_code=503, detail=f"Database unavailable: {exc}")

    # First wipe all existing embeddings for the folder so the reindex is clean
    file_ids = [f["id"] for f in resp.data or []]
    for fid in file_ids:
        try:
            supabase.table("embeddings").delete().eq("file_id", fid).execute()
        except Exception:
            pass

    for f in resp.data or []:
        try:
            index_file(folder_id, f["id"], f["storage_path"], f["name"])
        except Exception as exc:
            print(f"[index_folder] ERROR indexing file {f['name']!r}: {exc}", flush=True)
            continue   # don't abort the whole reindex if one file fails


# -------------------- DELETE EMBEDDINGS (helper for main.py) --------------------

def delete_embeddings_for_folder(folder_id: str):
    """Delete all embeddings rows for a folder (called before reindex or delete)."""
    try:
        files_resp = supabase.table("files").select("id").eq("folder_id", folder_id).execute()
        file_ids = [f["id"] for f in files_resp.data or []]
        if file_ids:
            supabase.table("embeddings").delete().in_("file_id", file_ids).execute()
    except Exception:
        pass


def delete_embeddings_for_file(file_id: str):
    """Delete all embeddings rows for a specific file."""
    try:
        supabase.table("embeddings").delete().eq("file_id", file_id).execute()
    except Exception:
        pass


# -------------------- FOLDER SELECTION --------------------

def find_best_folder(question: str) -> str | None:
    """
    Query every folder via pgvector similarity and return the one whose
    best chunk has the highest cosine similarity to `question`.
    Returns None if no folder has any embeddings.
    """
    try:
        folders_resp = execute_with_retry(
            lambda sb: sb.table("folders").select("id, name").execute()
        )
    except Exception as exc:
        raise HTTPException(status_code=503, detail=f"Database unavailable: {exc}")

    folders = folders_resp.data or []
    if not folders:
        return None

    q_vec = embed_texts([question])[0].tolist()

    best_folder_name = None
    best_similarity  = -1.0

    for folder in folders:
        try:
            result = supabase.rpc(
                "match_chunks",
                {
                    "query_embedding":   q_vec,
                    "target_folder_id":  folder["id"],
                    "match_count":       1,
                }
            ).execute()
        except Exception:
            continue

        rows = result.data or []
        if rows:
            sim = float(rows[0].get("similarity", 0))
            if sim > best_similarity:
                best_similarity  = sim
                best_folder_name = folder["name"]

    return best_folder_name

_SUMMARIZE_KEYWORDS = {
    "summarize", "summary", "summarise", "overview", "brief",
    "outline", "describe", "what is this", "what does this",
    "what is the document", "explain this", "tldr", "tl;dr",
}

def _is_summarize_intent(question: str) -> bool:
    q = question.lower()
    return any(kw in q for kw in _SUMMARIZE_KEYWORDS)


def _fetch_all_chunks(folder_id: str, file_name: str | None = None) -> list[str]:
    """Fetch chunks for a folder. If file_name is given, only fetch that file's chunks."""
    try:
        files_q = supabase.table("files").select("id, name").eq("folder_id", folder_id)
        if file_name:
            files_q = files_q.ilike("name", file_name)
        files_resp = files_q.execute()
        file_ids = [f["id"] for f in files_resp.data or []]
        if not file_ids:
            return []
        chunks_resp = (
            supabase.table("embeddings")
            .select("chunk_text")
            .in_("file_id", file_ids)
            .execute()
        )
        return [r["chunk_text"] for r in chunks_resp.data or []]
    except Exception:
        return []


# -------------------- RAG TOOL ENTRY --------------------

def run(question: str, folder: str = None):
    """
    Run the RAG pipeline.

    - folder: optional. If omitted, find_best_folder() automatically
      picks the most relevant project by vector similarity.
    - Summarization intent is detected automatically and uses all chunks
      instead of top-K similarity search.
    """
    # ── Resolve folder ──────────────────────────────────────
    if not folder:
        folder = find_best_folder(question)

    if not folder:
        return {"error": "No documents found in any project folder"}

    # ── Fetch folder record ─────────────────────────────────
    try:
        resp = execute_with_retry(
            lambda sb: (
                sb.table("folders")
                .select("id")
                .eq("name", folder)
                .single()
                .execute()
            )
        )
    except Exception as exc:
        raise HTTPException(status_code=503, detail=f"Database unavailable: {exc}")

    if not resp.data:
        return {"error": "Folder not found"}

    folder_id = resp.data["id"]

    # ── Choose retrieval strategy ───────────────────────────
    if _is_summarize_intent(question):
        # Check if a specific filename is mentioned in the question
        mentioned_file = None
        try:
            files_resp = supabase.table("files").select("name").eq("folder_id", folder_id).execute()
            for f in files_resp.data or []:
                if f["name"].lower() in question.lower():
                    mentioned_file = f["name"]
                    break
        except Exception:
            pass

        # Summarization: use chunks from the specific file or all files
        all_chunks = _fetch_all_chunks(folder_id, mentioned_file)
        if not all_chunks:
            return {"error": "No indexed content found in this folder. Please click \"Re-index\" to index the uploaded files before querying."}

        # Truncate to avoid LLM context limits (~6000 words)
        context = "\n\n".join(all_chunks)
        words = context.split()
        if len(words) > 6000:
            context = " ".join(words[:6000]) + "\n\n[... document truncated for length ...]"

        prompt = f"""You are a helpful assistant.
Provide a clear, structured summary of the following document content.
Include the main topics, key points, and any important details.

Document Content:
{context}

Task: {question}
"""
    else:
        # Normal RAG: top-K vector similarity search
        q_vec = embed_texts([question])[0].tolist()

        try:
            result = supabase.rpc(
                "match_chunks",
                {
                    "query_embedding":   q_vec,
                    "target_folder_id":  folder_id,
                    "match_count":       TOP_K,
                }
            ).execute()
        except Exception as exc:
            raise HTTPException(status_code=503, detail=f"Database unavailable: {exc}")

        rows = result.data or []
        if not rows:
            return {"error": "No indexed content found in this folder. Please click \"Re-index\" to index the uploaded files before querying."}

        context = "\n\n".join(r["chunk_text"] for r in rows)

        prompt = f"""You are a helpful assistant.
Answer the question ONLY using the context below.
If the answer is not found in the context, say:
"I don't know based on the provided documents."

Context:
{context}

Question:
{question}
"""

    answer = generate_response(prompt)

    sources = []
    if _is_summarize_intent(question):
        sources = [{"label": mentioned_file or "All Folder Documents", "score": 1.0}]
    else:
        for i, r in enumerate(rows):
            sim = float(r.get("similarity", 0))
            if sim > 0:
                sources.append({"label": f"Knowledge Match #{i+1}", "score": round(sim, 2)})

    return {
        "type":        "document",
        "folder_used": folder,
        "answer":      answer.strip(),
        "sources":     sources
    }